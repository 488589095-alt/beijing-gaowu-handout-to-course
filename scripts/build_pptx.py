# -*- coding: utf-8 -*-
"""
build_pptx.py — 北京高中物理「content.json → 课件 pptx」渲染引擎

无独立 PPT 模版：直接打开【参考PPT】，复用它的「具名版式」(首页/章节页/知识讲解/典型例题/
真题呈现/知识总结页/结束) —— 这些版式自带品牌背景与页眉装饰；在其上铺讲义内容，
最后删掉参考PPT原有的成品页，只留生成页。

物理三件核心能力：
  · 公式：每个 <m:oMath> 渲染成「a14 原生(PowerPoint可编辑) + 图片兜底(全阅读器/本机可见)」
    的独立形状（见 omml.math_shape_xml）；简单行内公式走 Unicode 文本直接进句子。
  · 图：光路图随锚点内联放到所在知识块/例题页。
  · 题：讲义无答案 → 例题/真题「题页 → 答案页(AI生成·待师审)」两页式，答案标红。

用法:
  python3 build_pptx.py --content output/<讲次>/content.json \
      --ref "<参考PPT.pptx>" -o output/<讲次>/<讲次>_课件.pptx
"""
import argparse, copy, json, re, sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from lxml import etree
from PIL import Image

sys.path.insert(0, str(Path(__file__).parent))
from omml import (math_shape_xml, math_native_xml, math_inline_xml,  # noqa: E402
                  render_math_images)

SHAPE_TAGS = ("sp", "pic", "grpSp", "graphicFrame", "cxnSp")
REF_ATTRS = [qn("r:embed"), qn("r:link"), qn("r:id"), qn("r:pict")]
FORMULA_MODE = "native"   # native=纯公式格式(WPS/PPT可编辑) | hybrid=公式+图片兜底(全阅读器可见)


def clone_slide(prs, src):
    """深拷贝参考PPT的某页所有形状到新页（同版式），保留图形/字体/位置/图片关系。"""
    new = prs.slides.add_slide(src.slide_layout)
    spTree = new.shapes._spTree
    for ch in list(spTree):
        if etree.QName(ch).localname in SHAPE_TAGS:
            spTree.remove(ch)
    for el in src.shapes._spTree:
        if etree.QName(el).localname in SHAPE_TAGS:
            spTree.append(copy.deepcopy(el))
    relmap = {}
    for rId, rel in src.part.rels.items():
        if rel.reltype.endswith("slideLayout"):
            continue
        new_r = (new.part.rels.get_or_add_ext_rel(rel.reltype, rel._target)
                 if rel.is_external else new.part.relate_to(rel.target_part, rel.reltype))
        relmap[rId] = new_r
    for el in spTree.iter():
        for a in REF_ATTRS:
            if el.get(a) in relmap:
                el.set(a, relmap[el.get(a)])
    return new


def replace_text(shape, new_text):
    """换文字，保留首 run 的全部样式（字体/字号/色/位置），删除其余 run/段。"""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    runs = list(p0.runs)
    if not runs:
        _set_run(p0.add_run(), new_text, None, None, None, ea=None, latin=None)
        return
    r0 = runs[0]
    r0.text = new_text
    for r in runs[1:]:
        r._r.getparent().remove(r._r)
    for extra in tf.paragraphs[1:]:
        extra._p.getparent().remove(extra._p)


def find_fixtures(prs):
    """从参考PPT原页里识别可克隆的 fixture：封面/模块标题/结束。"""
    fx = {"cover": None, "module": None, "end": None}
    for s in prs.slides:
        ln = s.slide_layout.name
        txt = " ".join(sh.text_frame.text for sh in s.shapes
                       if sh.has_text_frame)
        if fx["cover"] is None and ln in ("首页", "课节首页", "封面"):
            fx["cover"] = s
        if fx["module"] is None and ("模块" in txt or ln in ("章节下的小标题", "模块标题", "章节页")):
            if any(sh.has_text_frame and any(r.font.size and r.font.size.pt >= 48
                   for p in sh.text_frame.paragraphs for r in p.runs)
                   for sh in s.shapes):
                fx["module"] = s
        if ("本节课结束" in txt or "再见" in txt):
            fx["end"] = s
    if fx["cover"] is None:
        fx["cover"] = prs.slides[0]
    return fx

# ════════ 设计 token（从参考PPT逐run拆解：references/style_form.md）════════
INK = "000000"
RED = "FF0000"          # 强调/答案（参考PPT正文强调红）
RED2 = "C00000"         # Part 标签深红
AI_PURPLE = "7A2E8E"    # AI 生成注释
EA = "微软雅黑"          # 中文正文/标题
LATIN = "Times New Roman"  # 西文/数字
EA_TITLE = "微软雅黑"
PAGE_W, PAGE_H = 10.0, 7.5
AI_NOTE = "（AI 生成 · 待老师审核修改）"
CIRC = "①②③④⑤⑥⑦⑧⑨⑩⑪⑫⑬⑭"

# 内容区（知识讲解/典型例题 版式：页眉占 0~1.15，内容 1.25~7.15）
BODY_X, BODY_W = 0.5, 9.0
BODY_TOP, BODY_BOT = 1.3, 7.15
TITLE_POS = (0.6, 0.22, 7.0, 0.6)   # 页内小标题槽（参考"光的折射率"sz21.6 B）

MEDIA_DIR = None   # render() 内设定
MATH_IMG = {}      # mid -> png path（块公式兜底图）
_SID = [1000]


def _sid():
    _SID[0] += 1
    return _SID[0]


# ════════ 版式查找（按角色给多候选名，兼容不同参考PPT命名）════════
# 同一角色在不同讲次参考PPT里命名不同（test1 首页/章节页/知识讲解/典型例题；
# test2 课节首页/模块标题/知识精讲/例题精讲）→ 每个角色列全部候选，命中即用。
ROLE_LAYOUTS = {
    "cover":   ["首页", "课节首页", "封面"],
    "divider": ["章节页", "【模块】模块封面", "模块标题", "模块封面"],
    "section": ["章节下的小标题", "模块标题", "章节页"],
    "knowledge": ["知识讲解", "【讲解】知识精讲", "知识精讲", "内容页"],
    "example": ["典型例题", "【讲解】例题精讲", "例题精讲", "知识讲解"],
    "exam":    ["真题呈现", "【讲解】例题精讲", "典型例题", "例题精讲"],
    "summary": ["知识总结页", "小结页——知识", "小结页——题型", "知识讲解"],
    "end":     ["4_自定义版式", "未来你会遇见首页", "首页", "课节首页"],
    "blank":   ["内容页", "知识讲解", "知识精讲"],
}


def layout_for(prs, role):
    for nm in ROLE_LAYOUTS.get(role, []):
        for L in prs.slide_layouts:
            if L.name == nm:
                return L
    return prs.slide_layouts[0]


def strip_placeholders(slide):
    """删掉从版式继承来的空占位符（否则显示『单击此处添加文本』/图片图标占位）。"""
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)


def new_slide(prs, role):
    s = prs.slides.add_slide(layout_for(prs, role))
    strip_placeholders(s)   # 清空占位符，只保留版式自带的品牌/页眉装饰
    return s


# ════════ 文本 ════════
def _set_run(r, text, sz, bold, color, ea=EA, latin=LATIN):
    r.text = text
    if sz is not None:
        r.font.size = Pt(sz)
    if bold is not None:
        r.font.bold = bold
    if color is not None:
        r.font.color.rgb = RGBColor.from_string(color)
    rpr = r._r.get_or_add_rPr()
    for tag, v in (("ea", ea), ("latin", latin)):
        if v:
            e = rpr.find(qn("a:" + tag))
            if e is None:
                e = etree.SubElement(rpr, qn("a:" + tag))
            e.set("typeface", v)


_GREEK = set("αβγδεζηθικλμνξοπρστυφχψωΓΔΘΛΞΠΣΦΨΩ")
_SUBS = set("₀₁₂₃₄₅₆₇₈₉₊₋₌₍₎ₐₑₒₓₕₖₗₘₙₚₛₜⁿ⁰¹²³⁴⁵⁶⁷⁸⁹")


def add_var_runs(p, text, sz, color=INK, bold=False):
    """把含变量的数学文本逐段成 run：单个拉丁/希腊字母=斜体(变量 n,c,θ)，
       多字母(sin/cm 等函数/单位)、数字、下标、运算符=正体。符合物理排版惯例。"""
    i, n = 0, len(text)
    while i < n:
        c = text[i]
        if c.isalpha() and c.isascii():       # 拉丁字母连续段
            j = i
            while j < n and text[j].isascii() and text[j].isalpha():
                j += 1
            seg = text[i:j]
            _set_run(p.add_run(), seg, sz, bold, color,
                     ea=None, latin="Times New Roman")
            if len(seg) == 1:                  # 单字母=变量→斜体
                p.runs[-1].font.italic = True
            i = j
        elif c in _GREEK:                      # 希腊字母=变量→斜体
            _set_run(p.add_run(), c, sz, bold, color, ea=None, latin=None)
            p.runs[-1].font.italic = True
            i += 1
        else:                                  # 数字/下标/运算符/单位字符=正体
            j = i
            while j < n and not (text[j].isascii() and text[j].isalpha()) \
                    and text[j] not in _GREEK:
                j += 1
            _set_run(p.add_run(), text[i:j], sz, bold, color, ea=DEFAULT_EA_FOR(text[i:j]))
            i = j


def DEFAULT_EA_FOR(s):
    return EA if any("一" <= c <= "鿿" for c in s) else None


def fill_para(p, segs, sz, color=INK, bold=False, math_pool=None, drop_label=False):
    """统一把 segments 填进段落 p：文字→run；行内简单公式→Unicode变量斜体；
       块公式→**行内原生公式**(嵌进本段,不另起形状·不偏移)。返回是否含块公式。"""
    # 给段落写入 defRPr sz，让行内 <a14:m> 公式区继承与正文相同的字号
    pPr = p._p.get_or_add_pPr()
    defRPr = pPr.find(qn("a:defRPr"))
    if defRPr is None:
        defRPr = etree.SubElement(pPr, qn("a:defRPr"))
    defRPr.set("sz", str(int(sz * 100)))
    has_block = False
    for s in segs:
        if s["type"] == "text":
            t = s["text"].replace("\t", "").strip("\n")
            if drop_label:
                t = re.sub(r"^\s*[A-E][．.。]\s*", "", t); drop_label = False
            if t:
                _set_run(p.add_run(), t, sz, bold, color)
        elif s["type"] == "math":
            if s.get("inline"):
                u = s.get("unicode", "")
                if u:
                    add_var_runs(p, u, sz, color, bold)
            elif math_pool and s["id"] in math_pool:
                om = etree.fromstring(math_pool[s["id"]])
                p._p.append(etree.fromstring(math_inline_xml(om)))
                has_block = True
    return has_block


def textbox(slide, x, y, w, h, anchor=None, wrap=True):
    bx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = bx.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE   # 关 spAutoFit：框不自动长高去压下一框(PATCH-001)
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    if anchor is not None:
        tf.vertical_anchor = anchor
    return bx, tf


def title_box(slide, text, sz=22):
    bx, tf = textbox(slide, *TITLE_POS, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    _set_run(p.add_run(), text, sz, True, INK)
    return bx


def ai_note(slide):
    bx, tf = textbox(slide, 0.5, PAGE_H - 0.55, PAGE_W - 1, 0.4)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run(), AI_NOTE, 13, False, AI_PURPLE)


DIFFICULTY = "基础"   # render() 从讲义标题【基础/提升/拔高】解析；旗帜数=难度档
_DIFF_FLAGS = {"基础": 1, "提升": 2, "拔高": 3, "高考": 3}


FLAG_PNG = None   # render() 从参考PPT画布外抽出的真实难度旗帜图


def diff_tag(slide, page=""):
    """例题/真题页右上角标（仿参考PPT）：讲义页码圆角蓝框(P__待老师补) + 难度旗帜(靠右边·框外)。
       旗帜优先用参考PPT原图(FLAG_PNG)，无则文字▶。页码讲义docx无来源→占位 P__ 待老师补。"""
    from pptx.enum.shapes import MSO_SHAPE
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(7.5), Inches(0.22), Inches(1.3), Inches(0.4))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor.from_string("E8F4FF")
    box.line.color.rgb = RGBColor.from_string("0EA5FF"); box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = False
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run(), page or "P__", 16, True, "0EA5FF")
    # 难度旗帜（靠页面右上·框外）
    if FLAG_PNG and Path(FLAG_PNG).exists():
        fw, fh = fit_size(FLAG_PNG, 1.1, 0.4)
        slide.shapes.add_picture(FLAG_PNG, Inches(8.85), Inches(0.2),
                                 Inches(fw), Inches(fh))
    else:
        n = _DIFF_FLAGS.get(DIFFICULTY, 1)
        fb = slide.shapes.add_textbox(Inches(8.9), Inches(0.2), Inches(1.05), Inches(0.42))
        _set_run(fb.text_frame.paragraphs[0].add_run(), "▶" * n, 18, True, "FE6C29")
    return box


def harvest_flag(prs, out_dir):
    """从参考PPT例题页画布外(x<0或x>10)抽出难度旗帜图，按难度档挑一张存盘，返回路径。"""
    from pptx.util import Emu
    cands = {}
    for s in prs.slides:
        for sh in s.shapes:
            if sh.shape_type == 13 and sh.left is not None:
                lx = Emu(sh.left).inches
                if lx < -0.3 or lx > 10.2:
                    try:
                        w = Emu(sh.width).inches
                        cands[round(w, 2)] = sh.image
                    except Exception:
                        pass
    if not cands:
        return None
    widths = sorted(cands)                      # 越宽=旗帜越多
    idx = min(_DIFF_FLAGS.get(DIFFICULTY, 1) - 1, len(widths) - 1)
    # 基础→偏少旗帜；按难度档在已抽到的旗帜里选一张（窄→宽）
    pick = cands[widths[max(0, min(idx, len(widths) - 1))]]
    out = Path(out_dir) / f"_flag.{pick.ext}"
    out.write_bytes(pick.blob)
    return str(out)


# ════════ 图片 ════════
def _aspect(png):
    try:
        im = Image.open(png); return im.size[0] / im.size[1] if im.size[1] else 1
    except Exception:
        return 1


def fit_size(png, max_w, max_h):
    try:
        im = Image.open(png); iw, ih = im.size
    except Exception:
        return max_w, max_h
    ar = iw / ih if ih else 1
    w, h = max_w, max_w / ar
    if h > max_h:
        h = max_h; w = max_h * ar
    return w, h


def add_image(slide, png, x, y, max_w, max_h, center_x=True):
    if not png or not Path(png).exists():
        return y
    w, h = fit_size(png, max_w, max_h)
    px = x + (max_w - w) / 2 if center_x else x
    slide.shapes.add_picture(png, Inches(px), Inches(y), Inches(w), Inches(h))
    return y + h


# ════════ 公式（a14 原生 + 图片兜底）独立形状 ════════
def add_math(slide, mid, om_xml, x, y, w, h, sz=18):
    """放一个块公式形状。native=纯原生公式(WPS/PPT可编辑,默认)；hybrid=原生+图片兜底。"""
    om = etree.fromstring(om_xml)
    if FORMULA_MODE == "native":
        xml = math_native_xml(om, _sid(), x, y, w, h, sz=sz)
        slide.shapes._spTree.append(etree.fromstring(xml))
        return y + h
    png = MATH_IMG.get(mid)
    rid = None
    if png and Path(png).exists():
        w, h = fit_size(png, w, h)
        pic = slide.shapes.add_picture(png, Inches(x), Inches(y), Inches(w), Inches(h))
        rid = pic._element.blip_rId
        pic._element.getparent().remove(pic._element)
    xml = math_shape_xml(om, _sid(), x, y, w, h, img_rel_id=rid)
    slide.shapes._spTree.append(etree.fromstring(xml))
    return y + h


# ════════ 流式排版引擎（自动分页）════════
class Flow:
    """在内容区按 段落/公式/图片 自上而下流式排版；溢出时回调 new_page() 取新页。"""
    def __init__(self, prs, new_page, x=BODY_X, w=BODY_W, top=BODY_TOP, bot=BODY_BOT):
        self.prs = prs
        self.new_page = new_page
        self.x, self.w, self.top, self.bot = x, w, top, bot
        self.slide = None
        self.y = top

    def ensure(self, need):
        if self.slide is None or self.y + need > self.bot:
            self.slide = self.new_page()
            self.y = self.top
        return self.slide

    @staticmethod
    def _char_w(c, sz):
        """字符视觉宽度(in)。中文/全角标点≈1.0em；西文/数字/半角≈0.55em。
           物理排版必须按全角算，否则中文行数严重低估→框过矮→叠印(PATCH-001)。"""
        em = sz / 72.0
        if ("一" <= c <= "鿿" or "　" <= c <= "〿"
                or "＀" <= c <= "￯"):
            return em
        return em * 0.55

    def _est_lines(self, text, sz, width=None):
        avail = (width or self.w) - 0.12
        total = 0
        for seg in text.split("\n"):
            if not seg:
                total += 1; continue
            w = 0.0; ln = 1
            for c in seg:
                cw = self._char_w(c, sz)
                if w + cw > avail:
                    ln += 1; w = cw
                else:
                    w += cw
            total += ln
        return max(1, total)

    def paragraph(self, segs, sz=16, bold=False, color=INK, label=None,
                  label_color=RED, gap=0.08, indent=0.0):
        """渲染一段（text+inline math unicode 混排）+ 段后块公式 + 段内图片。"""
        text_runs = []   # (text, is_math)
        block_maths = []
        imgs = []
        for s in segs:
            if s["type"] == "text":
                t = s["text"].replace("\t", "").strip("\n")
                if t:
                    text_runs.append((t, False))
            elif s["type"] == "math":
                if s.get("inline"):
                    u = s.get("unicode", "")
                    if u:
                        text_runs.append((u, True))
                else:
                    block_maths.append(s["id"])
            elif s["type"] == "img":
                imgs.append(s["ref"])

        full = (label or "") + "".join(t for t, _ in text_runs)
        if full.strip():
            lines = self._est_lines(full, sz, width=self.w - indent)
            lh = sz / 72 * 1.5
            need = lines * lh + gap
            self.ensure(need)
            bx, tf = textbox(self.slide, self.x + indent, self.y,
                             self.w - indent, need)
            p = tf.paragraphs[0]
            if label:
                _set_run(p.add_run(), label, sz, True, label_color)
            for t, is_math in text_runs:
                if is_math:
                    add_var_runs(p, t, sz, color, bold)   # 变量斜体/单位正体
                else:
                    _set_run(p.add_run(), t, sz, bold, color)
            self.y += need
        # 段后块公式（居中）
        for mid in block_maths:
            self._block_math(mid, sz)
        # 段内/段后图片
        for ref in imgs:
            self.image(ref)

    def _block_math(self, mid, sz=16):
        om_xml = self.math_pool.get(mid)
        if not om_xml:
            return
        png = MATH_IMG.get(mid)
        mh = 0.72   # 18pt 公式约占 0.25in/行，留 2~3 行余量
        mw = 3.5
        if png and Path(png).exists():
            mw, mh = fit_size(png, 5.0, 1.1)
        self.ensure(mh + 0.12)
        add_math(self.slide, mid, om_xml, self.x + 0.4, self.y, mw, mh, sz=sz)
        self.y += mh + 0.12

    def image(self, ref, max_w=5.5, max_h=3.0):
        png = str(MEDIA_DIR / ref) if MEDIA_DIR else None
        if not png or not Path(png).exists():
            return
        w, h = fit_size(png, max_w, max_h)
        self.ensure(h + 0.15)
        add_image(self.slide, png, self.x, self.y, self.w, h)
        self.y += h + 0.15

    def option(self, o, sz=15):
        """选项紧凑渲染：『A．』标签 + 同行接 文字/行内公式/块公式图；图形选项图放标签下。"""
        label = o.get("label", "")
        pieces = []   # (text, is_math)
        block_ids, imgs = [], []
        for s in o["segs"]:
            if s["type"] == "text":
                t = s["text"].replace("\t", "").strip()
                t = re.sub(r"^[A-E][．.。]\s*", "", t)   # 去掉重复的选项字母
                if t:
                    pieces.append((t, False))
            elif s["type"] == "math":
                if s.get("inline"):
                    u = s.get("unicode", "")
                    if u:
                        pieces.append((u, True))
                else:
                    block_ids.append(s["id"])
            elif s["type"] == "img":
                imgs.append(s["ref"])
        text = "".join(t for t, _ in pieces)
        lh = sz / 72 * 1.5
        self.ensure(lh)
        # 标签 + 行内文字（标签框窄，仅容 "A．"，不伸到正文下）
        bx, tf = textbox(self.slide, self.x, self.y, 0.5, lh, anchor=MSO_ANCHOR.MIDDLE)
        _set_run(tf.paragraphs[0].add_run(), f"{label}．", sz, True, INK)
        xcur = self.x + 0.55
        if text:
            bx, tf = textbox(self.slide, xcur, self.y, self.w - 0.6, lh,
                             anchor=MSO_ANCHOR.MIDDLE)
            p = tf.paragraphs[0]
            for t, is_math in pieces:
                if is_math:
                    add_var_runs(p, t, sz, INK)
                else:
                    _set_run(p.add_run(), t, sz, False, INK)
        # 块公式 紧跟标签同行（native 公式较高→预留空间防与下一选项相碰）
        rowh = lh
        for mid in block_ids:
            png = MATH_IMG.get(mid)
            mw, mh = (fit_size(png, 3.2, 0.72) if png and Path(png).exists()
                      else (2.0, 0.65))
            self.ensure(mh)
            add_math(self.slide, mid, self.math_pool[mid], xcur, self.y, mw, mh, sz=sz)
            xcur += mw + 0.2
            rowh = max(rowh, mh)
        self.y += rowh + 0.04
        # 图形选项（每个选项一张图）→ 放标签下方一行
        for ref in imgs:
            self.image(ref, max_w=4.0, max_h=2.2)
        self.y += 0.06


# ════════ 页渲染 ════════
def _maxfont(shape):
    return max((r.font.size.pt for p in shape.text_frame.paragraphs
               for r in p.runs if r.font.size), default=0)


def r_cover(prs, fx, C):
    """克隆参考PPT封面页，按内容启发式换文字（保留圆角框/吉祥物/精确字体字号）。"""
    s = clone_slide(prs, fx["cover"])
    texts = [sh for sh in s.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
    title_sh = max(texts, key=_maxfont, default=None)
    for sh in texts:
        t = sh.text_frame.text.strip()
        if re.match(r"第[一二三四五六七八九十百\d]+讲", t):
            replace_text(sh, C["lecture_no"])
        elif re.match(r"^高[一二三]$", t):
            replace_text(sh, C.get("grade", t))
        elif "老师" in t or "主讲" in t:
            replace_text(sh, C.get("teacher", "主讲老师："))
        elif sh is title_sh:
            replace_text(sh, C["title"])
    return s


def r_module_title(prs, fx, idx, name):
    """克隆参考PPT模块标题页（大字考点名 + 模块N），换文字。无fixture则退回版式。"""
    cn = ["一", "二", "三", "四", "五", "六", "七", "八"]
    if fx.get("module") is not None:
        s = clone_slide(prs, fx["module"])
        texts = [sh for sh in s.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
        big = max(texts, key=_maxfont, default=None)
        for sh in texts:
            t = sh.text_frame.text.strip()
            if re.match(r"^模块[一二三四五六七八九十]", t):
                replace_text(sh, f"模块{cn[idx-1]}")
            elif sh is big:
                replace_text(sh, name)
        return s
    s = new_slide(prs, "divider")
    bx, tf = textbox(s, 1.0, 2.6, 8.0, 1.9, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run(), name, 48, True, INK, EA_TITLE)
    bx, tf = textbox(s, 4.0, 1.7, 2.0, 0.7)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run(), f"模块{cn[idx-1]}", 26, False, INK)
    return s


def r_textbook(prs, C):
    if not C.get("textbook_link"):
        return
    s = new_slide(prs, "knowledge")
    title_box(s, "教材链接")
    rows = C["textbook_link"]
    add_table(s, rows, 0.6, 1.5, 8.8, min(5.0, 0.5 * len(rows)))


def add_table(slide, rows, x, y, w, h):
    nr, nc = len(rows), max(len(r) for r in rows)
    gf = slide.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(h))
    t = gf.table
    for ri, row in enumerate(rows):
        for ci in range(nc):
            cell = t.cell(ri, ci)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            txt = row[ci] if ci < len(row) else ""
            _set_run(p.add_run(), str(txt), 13 if ri else 14, ri == 0, INK)
    return gf


def r_knowledge(prs, sec, math_pool):
    """知识精讲：流式排版 sec['knowledge'] 各块，自动分页。"""
    state = {"first": True}

    def new_page():
        s = new_slide(prs, "knowledge")
        title_box(s, sec["name"])
        return s

    flow = Flow(prs, new_page)
    flow.math_pool = math_pool
    for blk in sec["knowledge"]:
        if blk["type"] == "table":
            s = flow.ensure(0.5 * len(blk["rows"]) + 0.2)
            add_table(s, blk["rows"], flow.x, flow.y, flow.w,
                      min(4.0, 0.45 * len(blk["rows"])))
            flow.y += min(4.0, 0.45 * len(blk["rows"])) + 0.15
            continue
        segs = blk["segs"]
        txt = "".join(s.get("text", "") for s in segs if s["type"] == "text").strip()
        # 概念小标题：(1)xxx / ①xxx / N.基础知识 → 加粗
        label = None
        import re as _re
        if _re.match(r"^\s*[（(]\d+[)）]", txt) or _re.match(r"^\s*\d+[．.]", txt) \
                or _re.match(r"^\s*[①-⑩]", txt):
            label = None  # 整段会以粗体呈现
            flow.paragraph(segs, sz=16, bold=True, color=INK, gap=0.06)
        else:
            flow.paragraph(segs, sz=15, bold=False, color=INK, gap=0.06)
        for ref in blk.get("imgs", []):
            flow.image(ref)


def _seg_paragraph_into(flow, segs, sz, color=INK, bold=False, label=None,
                        label_color=RED):
    flow.paragraph(segs, sz=sz, bold=bold, color=color, label=label,
                   label_color=label_color)


def _stem_lines(text, sz, w):
    avail = w - 0.12
    total = 0
    for seg in text.split("\n"):
        wd = 0.0; ln = 1
        for c in seg:
            cw = Flow._char_w(c, sz)
            if wd + cw > avail:
                ln += 1; wd = cw
            else:
                wd += cw
        total += ln
    return max(1, total)


def _para_text(segs):
    return "".join(s.get("text", "") if s["type"] == "text"
                   else (s.get("unicode", "") if s.get("inline") else "")
                   for s in segs)


def _render_stem(slide, segs, x, y, w, sz, math_pool=None, drop_label=False):
    """题干/正文：text + 行内公式(变量斜体) + 块公式(行内原生·不偏移) 混排；返回高度。"""
    full = _para_text(segs)
    n_block = sum(1 for s in segs if s["type"] == "math" and not s.get("inline"))
    lh = sz / 72 * 1.5
    h = _stem_lines(full, sz, w) * lh + n_block * (sz / 72 * 1.1) + 0.08
    bx, tf = textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP)
    p = tf.paragraphs[0]; p.line_spacing = 1.0
    fill_para(p, segs, sz, INK, math_pool=math_pool, drop_label=drop_label)
    return h, []


def _split_options(opts):
    """拆分被合并的选项：一段里含 'A．… B．…' → 拆成独立 A、B（讲义把两项写在一行）。"""
    out = []
    for o in opts:
        cur = {"label": o.get("label", ""), "segs": [], "imgs": list(o.get("imgs", []))}
        started = False
        for s in o["segs"]:
            if s["type"] == "text":
                parts = re.split(r"(?=[A-E][．.。])", s["text"])
                for part in parts:
                    mlab = re.match(r"\s*([A-E])[．.。]\s*(.*)", part, re.S)
                    if mlab:
                        if cur["segs"] or started:
                            out.append(cur)
                        cur = {"label": mlab.group(1), "segs": [], "imgs": []}
                        started = True
                        if mlab.group(2).strip():
                            cur["segs"].append({"type": "text", "text": mlab.group(2)})
                    elif part.strip():
                        cur["segs"].append({"type": "text", "text": part})
            elif s["type"] == "img":
                cur["imgs"].append(s["ref"]); started = True   # 图形选项的图→imgs(别丢)
            else:
                cur["segs"].append(s); started = True
        out.append(cur)
    return out


def _place_options(slide, ex, math_pool, x, y, w, bot, sz=18):
    """选项紧接题干往下：文字/公式选项竖排(标签+公式同行)；图形选项横排一行。返回结束 y。"""
    opts = _split_options(ex.get("options", []))
    if not opts:
        return y
    img_opts = any(o.get("imgs") for o in opts)
    if img_opts:
        n = len(opts)
        gap = 0.15
        cw = (w - gap * (n - 1)) / n
        ch = min(bot - y - 0.1, 2.0)
        for i, o in enumerate(opts):
            ox = x + i * (cw + gap)
            bx, tf = textbox(slide, ox, y, cw, 0.32)
            _set_run(tf.paragraphs[0].add_run(), f"{o['label']}．", sz, True, INK)
            for ref in o.get("imgs", []):
                png = str(MEDIA_DIR / ref) if MEDIA_DIR else None
                if png and Path(png).exists():
                    iw, ih = fit_size(png, cw, ch - 0.34)
                    add_image(slide, png, ox, y + 0.34, cw, ih)
        return y + ch + 0.1
    cy = y
    lh = sz / 72 * 1.5
    for o in opts:
        otext = _para_text(o["segs"])
        nlines = max(1, _stem_lines(otext, sz, w - 0.5))
        nblk = sum(1 for s in o["segs"] if s["type"] == "math" and not s.get("inline"))
        rowh = nlines * lh + nblk * (sz / 72 * 0.9) + 0.06
        bx, tf = textbox(slide, x, cy, w, rowh, anchor=MSO_ANCHOR.TOP)
        tf.word_wrap = True
        p = tf.paragraphs[0]; p.line_spacing = 1.0
        _set_run(p.add_run(), f"{o['label']}．", sz, True, INK)
        fill_para(p, o["segs"], sz, INK, math_pool=math_pool, drop_label=True)
        cy += rowh
    return cy


def _example_figs_and_leadin(ex):
    """从 stem_imgs + extra 里抽出 (图片refs 顺序, 引导文字seg列表)。"""
    figs, leadins = [], []
    for ref in ex.get("stem_imgs", []):
        figs.append(ref)
    for ext in ex.get("extra", []):
        for s in ext:
            if s["type"] == "img":
                figs.append(s["ref"])
        if any(s["type"] != "img" and (s.get("text", "").strip() or s["type"] == "math")
               for s in ext):
            leadins.append([s for s in ext if s["type"] != "img"])
    # 题干 segs 内嵌的图
    for s in ex.get("stem_segs", []):
        if s["type"] == "img" and s["ref"] not in figs:
            figs.append(s["ref"])
    seen = set(); figs = [f for f in figs if not (f in seen or seen.add(f))]
    return figs, leadins


EX_SZ = 16   # 例题正文统一字号（同级内容一致，不逐页缩放；贴合参考PPT例题视觉字号）


def _opt_block_height(opts, sz, w):
    lh = sz / 72 * 1.5
    if any(o.get("imgs") for o in opts):
        return 2.1
    h = 0.0
    for o in _split_options(opts):
        ot = _para_text(o["segs"])
        nblk = sum(1 for s in o["segs"] if s["type"] == "math" and not s.get("inline"))
        h += max(1, _stem_lines(ot, sz, w - 0.5)) * lh + nblk * (sz / 72 * 0.9) + 0.06
    return h + 0.05


def _is_multiselect(ex, scaffold):
    """多选判定：选项含 E(≥5项) 或 AI答案含≥2个字母。"""
    nopt = len(_split_options(ex.get("options", [])))
    if nopt >= 5:
        return True
    ans = ((scaffold or {}).get(ex["no"], {}) or {}).get("answer", "")
    return len(re.findall(r"[A-E]", str(ans))) >= 2


def _mark_multiselect(stem_segs):
    """在题号后插『（多选）』(仿参考PPT)。"""
    out = [dict(s) for s in stem_segs]
    for s in out:
        if s["type"] == "text" and re.search(r"例题\s*\d+", s["text"]):
            s["text"] = re.sub(r"(例题\s*\d+)", r"\1（多选）", s["text"], count=1)
            break
    return out


def r_example(prs, ex, scaffold, math_pool, sec_name):
    """单页例题：题干→题图→选项 紧凑相邻、全在一页；右上=讲义页码+难度旗帜(真图)；
       答案/解析放画布外右侧(框外·不出现在放映)。字号统一 EX_SZ(同级一致)，仅缩题图保一页。"""
    s = new_slide(prs, "example")
    pg = ex.get("page")
    diff_tag(s, page=f"P{pg}" if pg else "")  # 右上：讲义页码(docx页) + 难度旗帜

    x, w, top, bot = 0.5, 9.0, 1.25, 7.05
    sz = EX_SZ
    figs, leadins = _example_figs_and_leadin(ex)
    lead_text = "".join(sg.get("text", "") for ld in leadins for sg in ld
                        if sg["type"] == "text")
    stem_segs = _mark_multiselect(ex["stem_segs"]) if _is_multiselect(ex, scaffold) \
        else ex["stem_segs"]
    figpaths = [str(MEDIA_DIR / r) for r in figs
                if MEDIA_DIR and (MEDIA_DIR / r).exists()]
    has_lead = bool(lead_text.strip())
    opt_h = _opt_block_height(ex.get("options", []), sz, w)

    def draw_figs(yy, max_h):
        """图块：多图并排(宽图在左,仿参考表左图右)，单图整宽；返回新 y。"""
        if not figpaths:
            return yy
        paths = sorted(figpaths, key=lambda p: -(_aspect(p)))   # 宽(表)在左
        if len(paths) >= 2:
            gap = 0.2; cw = (w - gap * (len(paths) - 1)) / len(paths)
            mih = 0
            for i, png in enumerate(paths):
                iw, ih = fit_size(png, cw, max_h)
                add_image(s, png, x + i * (cw + gap) + (cw - iw) / 2, yy, iw, ih,
                          center_x=False)
                mih = max(mih, ih)
            return yy + mih + 0.08
        iw, ih = fit_size(paths[0], min(w, 6.8), max_h)
        add_image(s, paths[0], x, yy, w, ih)
        return yy + ih + 0.08

    y = top
    h, _ = _render_stem(s, stem_segs, x, y, w, sz, math_pool=math_pool); y += h
    if has_lead:
        # 仿参考PPT：题干→引导→选项→图(底部大图)。图是供选项参考的资料(表/示意图)
        for ld in leadins:
            if any(sg["type"] != "img" for sg in ld):
                lh2, _ = _render_stem(s, ld, x, y, w, sz, math_pool=math_pool); y += lh2
        y = _place_options(s, ex, math_pool, x, y + 0.05, w, bot, sz=sz)
        draw_figs(y + 0.1, max(1.2, bot - y - 0.15))
    else:
        # 图是题目核心(光路图)：题干→图→选项
        if figpaths:
            avail = max(1.0, bot - y - opt_h - 0.12)
            y = draw_figs(y, min(avail, 2.8))
        _place_options(s, ex, math_pool, x, y + 0.05, w, bot, sz=sz)
    # 答案/解析：画布外右侧（仿参考 idx11 off-canvas，放映看不到）
    _offcanvas_answer(s, ex, scaffold)
    return s


def _offcanvas_answer(slide, ex, scaffold):
    sc = (scaffold or {}).get(ex["no"], {}) if scaffold else {}
    ans = sc.get("answer", "（待补）")
    ana = sc.get("analysis", "（解析待补）")
    bx, tf = textbox(slide, 10.4, 1.0, 5.6, 5.5, anchor=MSO_ANCHOR.TOP)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _set_run(p.add_run(), f"【{ex['no']} 答案】", 18, True, RED)
    _set_run(p.add_run(), str(ans), 18, True, RED)
    p2 = tf.add_paragraph()
    _set_run(p2.add_run(), "【解析】", 15, True, INK)
    for line in str(ana).split("\n"):
        pp = tf.add_paragraph()
        _set_run(pp.add_run(), line, 14, False, INK)
    pn = tf.add_paragraph()
    _set_run(pn.add_run(), AI_NOTE, 12, False, AI_PURPLE)


def r_summary(prs, scaffold, C):
    sm = (scaffold or {}).get("summary")
    if not sm:
        return
    s = new_slide(prs, "summary")
    bx, tf = textbox(s, 1.3, 0.4, 5.0, 0.6)
    _set_run(tf.paragraphs[0].add_run(), "知识总结", 22, True, INK)
    bx, tf = textbox(s, 1.0, 1.6, 8.0, 5.0, anchor=MSO_ANCHOR.TOP)
    first = True
    for line in (sm if isinstance(sm, list) else str(sm).split("\n")):
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        _set_run(p.add_run(), str(line), 18, False, INK)
    ai_note(s)
    return s


def r_end(prs, fx):
    if fx.get("end") is not None:
        return clone_slide(prs, fx["end"])   # 原样克隆结束页（含吉祥物/文字）
    s = new_slide(prs, "end")
    bx, tf = textbox(s, 0.68, 1.44, 4.0, 0.85)
    _set_run(tf.paragraphs[0].add_run(), "本节课结束", 33, False, INK, EA_TITLE)
    bx, tf = textbox(s, 0.78, 2.5, 6.0, 2.5)
    _set_run(tf.paragraphs[0].add_run(), "下节课我们", 40, True, INK, EA_TITLE)
    p = tf.add_paragraph()
    _set_run(p.add_run(), "再见啦～", 40, True, INK, EA_TITLE)
    return s


def _coverage_check(C, cdir, out_path):
    """生成后完整性校对：讲义图/公式/例题 是否都进了成片。漏则醒目告警。"""
    import zipfile, os
    mediadir = cdir / "media"
    hb = {}
    if mediadir.exists():
        for n in os.listdir(mediadir):
            if n.startswith("image"):
                hb[(mediadir / n).read_bytes()] = n
    z = zipfile.ZipFile(out_path)
    placed = {hb[z.read(n)] for n in z.namelist()
              if n.startswith("ppt/media/") and z.read(n) in hb}
    miss_img = sorted(set(hb.values()) - placed)
    a14 = sum(z.read(n).decode("utf-8", "ignore").count("<a14:m")
              for n in z.namelist()
              if n.startswith("ppt/slides/slide") and n.endswith(".xml"))
    n_block = 0
    for m in C["modules"]:
        for sec in m["sections"]:
            for b in sec["knowledge"]:
                if b["type"] == "para":
                    n_block += sum(1 for s in b["segs"]
                                   if s["type"] == "math" and not s.get("inline"))
            for e in sec["examples"]:
                for grp in [e["stem_segs"]] + e.get("extra", []) + \
                        [o["segs"] for o in e.get("options", [])]:
                    n_block += sum(1 for s in grp
                                   if s["type"] == "math" and not s.get("inline"))
    print("──完整性校对──")
    print(f"  讲义图 {len(hb)} → 成片 {len(placed)}",
          ("✅" if not miss_img else f"❌ 漏 {miss_img}"))
    print(f"  块公式应≥{n_block} → 成片原生 a14 {a14}",
          ("✅" if a14 >= n_block else "❌ 偏少"))
    if miss_img:
        print(f"  ⚠️ 有图未进成片：{miss_img}（请检查对应例题/知识页）")


# ════════ 主流程 ════════
def render(content_path, ref_path, out_path, formula_mode="native"):
    global MEDIA_DIR, MATH_IMG, FORMULA_MODE, DIFFICULTY
    FORMULA_MODE = formula_mode
    C = json.loads(Path(content_path).read_text(encoding="utf-8"))
    cdir = Path(content_path).parent
    MEDIA_DIR = cdir / "media"
    math_pool = C["math_pool"]
    m = re.search(r"【(基础|提升|拔高|高考)", C.get("title", "") + C.get("lecture_no", ""))
    DIFFICULTY = m.group(1) if m else "基础"

    scp = cdir / "ai_scaffold.json"
    scaffold = json.loads(scp.read_text(encoding="utf-8")) if scp.exists() else {}
    if not scaffold:
        print("⚠️ 无 ai_scaffold.json：答案/解析/小结将留空待补（见 SKILL.md 第2步）")

    # 渲染所有「块公式」兜底图（行内简单公式走 Unicode，不需图）
    block_ids = set()
    def scan(segs):
        for s in segs:
            if s.get("type") == "math" and not s.get("inline"):
                block_ids.add(s["id"])
    for m in C["modules"]:
        for sec in m["sections"]:
            for blk in sec["knowledge"]:
                if blk.get("type") in ("para",):
                    scan(blk["segs"])
            for ex in sec["examples"]:
                scan(ex["stem_segs"])
                for o in ex.get("options", []):
                    scan(o["segs"])
                for extra in ex.get("extra", []):
                    scan(extra)
    for ex in C.get("exam", []):
        scan(ex["stem_segs"])
        for o in ex.get("options", []):
            scan(o["segs"])
    if block_ids and FORMULA_MODE == "hybrid":
        print(f"渲染块公式兜底图：{len(block_ids)} 个 …")
        pool = {mid: math_pool[mid] for mid in block_ids if mid in math_pool}
        MATH_IMG = render_math_images(pool, MEDIA_DIR)
        print(f"  完成 {len(MATH_IMG)} 张")
    else:
        print(f"公式模式=native：{len(block_ids)} 个块公式以原生『公式格式』插入（WPS/PPT可编辑）")

    prs = Presentation(str(ref_path))
    orig = len(prs.slides)
    global FLAG_PNG
    FLAG_PNG = harvest_flag(prs, MEDIA_DIR)     # 挪用参考PPT画布外的真实难度旗帜图
    print(f"难度旗帜：{'已挪用参考PPT原图' if FLAG_PNG else '用文字▶兜底'}")
    fx = find_fixtures(prs)        # 先识别参考PPT可克隆 fixture（封面/模块标题/结束）
    print(f"克隆 fixture：封面={'✓' if fx['cover'] else '×'} "
          f"模块标题={'✓' if fx['module'] else '×'} 结束={'✓' if fx['end'] else '×'}")

    structure = []

    def rec(role, summ):
        structure.append({"n": len(prs.slides) - orig, "role": role, "summary": summ})

    r_cover(prs, fx, C); rec("cover", C["title"])
    r_textbook(prs, C); rec("textbook", "教材链接")
    for mi, m in enumerate(C["modules"], 1):
        r_module_title(prs, fx, mi, m["name"]); rec("module", m["name"])
        for si, sec in enumerate(m["sections"], 1):
            if sec["knowledge"]:
                r_knowledge(prs, sec, math_pool); rec("knowledge", sec["name"])
            for ex in sec["examples"]:
                r_example(prs, ex, scaffold, math_pool, sec["name"])
                rec("example", ex["no"])
    if C.get("exam"):
        r_module_title(prs, fx, len(C["modules"]) + 1, "真题呈现"); rec("module", "真题呈现")
        for ex in C["exam"]:
            r_example(prs, ex, scaffold, math_pool, "真题呈现"); rec("exam", ex["no"])
    r_summary(prs, scaffold, C); rec("summary", "知识总结")
    r_end(prs, fx); rec("end", "结束")

    # 删原参考PPT成品页（保留版式/母版/品牌）
    for i in range(orig, 0, -1):
        el = list(prs.slides._sldIdLst)[i - 1]
        prs.part.drop_rel(el.get(qn("r:id")))
        prs.slides._sldIdLst.remove(el)

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    _coverage_check(C, cdir, out_path)
    (cdir / "slide_structure.json").write_text(
        json.dumps({"total": len(prs.slides), "slides": structure},
                   ensure_ascii=False, indent=1), encoding="utf-8")
    print(f"✅ 课件已生成：{out_path}（{len(prs.slides)} 页）")


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--content", required=True)
    ap.add_argument("--ref", required=True, help="参考PPT（作版式/品牌/fixture来源）")
    ap.add_argument("--formula", choices=["native", "hybrid"], default="native",
                    help="native=纯公式格式(WPS/PPT可编辑,默认) | hybrid=公式+图片兜底(全阅读器可见)")
    ap.add_argument("-o", "--output", required=True)
    a = ap.parse_args()
    render(a.content, a.ref, a.output, a.formula)
