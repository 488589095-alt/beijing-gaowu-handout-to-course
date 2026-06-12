# -*- coding: utf-8 -*-
"""
blueprint_benchmark.py — 把标杆PPT完整拆解成"可理解的蓝图文本"（全页·逐文本形状）
用于：以标杆为蓝图 → 对应讲义填充 / 讲义无则AI生成。

每页输出：版式 + 每个文本形状(位置/尺寸/首run字体字号色/全文截断) + 图片·表格计数。
用法: python3 blueprint_benchmark.py <标杆.pptx> -o blueprint_raw.txt
"""
import argparse
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def IN(v):
    try:
        return round(Emu(v).inches, 1)
    except Exception:
        return None


def first_run_style(tf):
    for p in tf.paragraphs:
        for r in p.runs:
            if not r.text.strip():
                continue
            st = []
            if r.font.size:
                st.append(f"sz{r.font.size.pt:g}")
            if r.font.bold:
                st.append("B")
            try:
                if r.font.color and r.font.color.type is not None:
                    st.append("#" + str(r.font.color.rgb))
            except Exception:
                pass
            rpr = r._r.find(f"{A}rPr")
            if rpr is not None:
                ea = rpr.find(f"{A}ea")
                if ea is not None:
                    st.append(ea.get("typeface"))
            return ",".join(st)
    return ""


def walk(shapes, out, depth=0):
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            walk(sh.shapes, out, depth + 1)
            continue
        if sh.has_text_frame and sh.text_frame.text.strip():
            txt = " ⏎ ".join(p.text for p in sh.text_frame.paragraphs).strip()
            pos = f"({IN(sh.left)},{IN(sh.top)} {IN(sh.width)}x{IN(sh.height)})"
            out.append(f"   ·{'  '*depth} {pos} [{first_run_style(sh.text_frame)}] 「{txt[:90]}」")
        elif sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            out.append(f"   ·{'  '*depth} IMG ({IN(sh.left)},{IN(sh.top)} {IN(sh.width)}x{IN(sh.height)})")
        elif sh.has_table:
            tb = sh.table
            cells = " | ".join(c.text.replace("\n", "/")[:18] for c in tb.rows[0].cells)
            out.append(f"   · TABLE {len(tb.rows)}x{len(tb.columns)} 行0:[{cells}]")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("-o", "--out", required=True)
    a = ap.parse_args()
    prs = Presentation(a.pptx)
    out = [f"# 蓝图: {Path(a.pptx).name}  共{len(prs.slides)}页  画布{IN(prs.slide_width)}x{IN(prs.slide_height)}in"]
    for i, s in enumerate(prs.slides, 1):
        out.append(f"\n=== S{i} [{s.slide_layout.name}] ===")
        walk(s.shapes, out)
    Path(a.out).write_text("\n".join(out), encoding="utf-8")
    print(f"→ {a.out} ({len(out)} 行)")


if __name__ == "__main__":
    main()
