# -*- coding: utf-8 -*-
"""
extract_styles.py — 参考PPT → 全局文字样式拆解（逐 shape 逐 run，混排不漏）
用于制作 references/style_form.json（生成器唯一样式来源）。

用法: python3 extract_styles.py <参考.pptx> --slides 1,3,4 [-o raw.json]
输出: 每页每文本shape的全部 run 样式 (sz/bold/color/latin/ea/align) + 文字片段
"""
import argparse, json
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def IN(v):
    try:
        return round(Emu(v).inches, 2)
    except Exception:
        return None


def run_style(r):
    st = {"snippet": r.text[:28]}
    if r.font.size:
        st["sz"] = r.font.size.pt
    if r.font.bold is not None:
        st["bold"] = r.font.bold
    try:
        if r.font.color and r.font.color.type is not None:
            st["color"] = str(r.font.color.rgb)
    except Exception:
        pass
    rpr = r._r.find(f"{A}rPr")
    if rpr is not None:
        for tag in ("latin", "ea"):
            e = rpr.find(f"{A}{tag}")
            if e is not None:
                st[tag] = e.get("typeface")
    return st


def walk(shapes, out):
    for sh in shapes:
        if sh.shape_type == 6:  # GROUP
            walk(sh.shapes, out)
            continue
        if not getattr(sh, "has_text_frame", False) or not sh.text_frame.text.strip():
            continue
        d = {"shape": sh.name, "pos": [IN(sh.left), IN(sh.top)],
             "size": [IN(sh.width), IN(sh.height)], "paras": []}
        for p in sh.text_frame.paragraphs:
            runs, last = [], None
            for r in p.runs:
                st = run_style(r)
                key = {k: v for k, v in st.items() if k != "snippet"}
                if runs and key == {k: v for k, v in runs[-1].items()
                                    if k not in ("snippet",)}:
                    runs[-1]["snippet"] = (runs[-1]["snippet"] + st["snippet"])[:28]
                else:
                    runs.append(st)
            if runs:
                d["paras"].append({"align": str(p.alignment) if p.alignment else None,
                                   "runs": runs})
        if d["paras"]:
            out.append(d)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--slides", required=True, help="逗号分隔页码(1-based)")
    ap.add_argument("-o", "--out", default=None)
    a = ap.parse_args()
    prs = Presentation(a.pptx)
    want = [int(x) for x in a.slides.split(",")]
    res = {}
    for no in want:
        if no <= len(prs.slides):
            out = []
            walk(prs.slides[no - 1].shapes, out)
            res[str(no)] = out
    txt = json.dumps(res, ensure_ascii=False, indent=1)
    if a.out:
        Path(a.out).write_text(txt, encoding="utf-8")
        print(f"→ {a.out}")
    else:
        print(txt)


if __name__ == "__main__":
    main()
